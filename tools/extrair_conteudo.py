"""Tira do PPTX oficial do mês o conteúdo que NÃO sai de planilha.

Comentários do DRE, programação e resultados de exposições, histórico de manejo
e fotos: nada disso tem base de dados por trás — é escrito e fotografado a cada
mês. Em vez de deixar esses slides como placeholder, este script lê o deck
aprovado daquele mês e grava o conteúdo em `_docs/comite_conteudo.json`, que
passa a ser o arquivo que a pessoa edita quando quiser ajustar algo.
`tools/build_comite.py` lê esse JSON.

Antes era fixo em junho/2026, com número de slide HARDCODED e sobrescrevia o
JSON inteiro — rodar de novo pra outro mês apagava os meses já semeados. Agora
acha os slides pelo TÍTULO (robusto a mês ganhar/perder um slide — julho tem 5
exposições com resultado, junho tinha 4) e faz MERGE no JSON, uma chave por mês.

A leitura de cada slide é por geometria: o gerador do deck original não usou
tabelas de verdade, só caixas de texto posicionadas, então cada slide é
reconstruído agrupando as caixas por linha (mesmo `top`) e ordenando por `left`.

Fotos: o deck oficial agrupa por TEMA ("Obras e melhorias realizadas · Banqueta",
um ou mais slides por tema) — não é uma grade solta do mês. `build_comite.py`
espera essa mesma forma (ver `fotos: [{"tema":..., "arquivos":[...]}]`).

Uso:
    python tools/extrair_conteudo.py <caminho.pptx> <AAAA-MM>
"""
from __future__ import annotations

import io
import json
import re
import sys
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
OUT = REPO / "_docs" / "comite_conteudo.json"
FOTOS = REPO / "assets" / "comite" / "fotos"

EMU_IN = 914400
LARG_FOTO = 1400                     # px — foto de slide não precisa de mais


def _titulo(slide) -> str:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip()
    return ""


def _acha(slides, contem) -> int | None:
    """Índice (0-based) do primeiro slide cujo título CONTÉM `contem`."""
    for i, sl in enumerate(slides):
        if contem.upper() in _titulo(sl).upper():
            return i
    return None


def _acha_todos(slides, contem) -> list[int]:
    return [i for i, sl in enumerate(slides) if contem.upper() in _titulo(sl).upper()]


def caixas(slide):
    """[(top_in, left_in, texto)] das caixas com texto, em ordem de leitura."""
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t:
            out.append(((sh.top or 0) / EMU_IN, (sh.left or 0) / EMU_IN, t))
    return sorted(out, key=lambda x: (round(x[0], 1), x[1]))


def linhas(slide, tol=0.12, desde=0.9):
    """Agrupa as caixas em linhas: mesma altura (± tol) = mesma linha."""
    grupos = []
    for top, left, txt in caixas(slide):
        if top < desde:
            continue
        if grupos and abs(grupos[-1][0] - top) <= tol:
            grupos[-1][1].append((left, txt))
        else:
            grupos.append((top, [(left, txt)]))
    return [[t for _, t in sorted(g, key=lambda x: x[0])] for _, g in grupos]


# --------------------------------------------------------- COMENTÁRIOS YTD
def comentarios(slide):
    """3 caixas por linha: categoria · texto · delta. O delta às vezes cai na
    linha de baixo (a caixa é mais alta), então casa pelo formato R$."""
    itens, pend = [], None
    for top, left, txt in caixas(slide):
        if top < 0.9:
            continue
        if re.fullmatch(r"[+\-−]?R\$ ?[\d.,]+k?", txt):
            if pend:
                pend["delta"] = txt
                itens.append(pend)
                pend = None
            continue
        if left < 1.5:
            if pend:
                itens.append(pend)
            pend = {"cat": txt, "txt": "", "delta": ""}
        elif pend:
            pend["txt"] = txt
    if pend:
        itens.append(pend)
    return [i for i in itens if i["txt"]]


# ------------------------------------------------------------ PROGRAMAÇÃO
def programacao(slide):
    ls = linhas(slide, desde=1.0)
    if not ls:
        return []
    return [l for l in ls[1:] if len(l) >= 3]      # ls[0] é o cabeçalho


# -------------------------------------------------------------- RESULTADOS
def resultados(slide):
    """Duas colunas de blocos: NOME DO ANIMAL em caixa alta, prêmios com 🏆."""
    cx = caixas(slide)
    titulo = cx[0][2] if cx else ""
    sub = next((t for _, l, t in cx if l < 1.0 and t != titulo), "")
    animais, atual = [], None
    for top, left, txt in sorted(cx, key=lambda x: (x[1] > 4.5, x[0])):
        if top < 0.9:
            continue
        if txt.startswith("🏆"):
            if atual:
                atual["premios"].append(txt.lstrip("🏆 ").strip())
            continue
        if txt == sub:
            continue
        atual = {"nome": txt, "premios": []}
        animais.append(atual)
    return {"titulo": titulo, "sub": sub, "animais": [a for a in animais if a["premios"]]}


# ------------------------------------------------------------------ MANEJO
def manejo(slide):
    """Mês curto numa caixa à esquerda, texto do lado.

    Os dois não têm o mesmo `top` — a caixa do mês fica alguns centésimos ABAIXO
    da do texto. Percorrer em ordem de leitura casava o texto com o mês seguinte
    (Jan recebia o texto de Fev) e perdia o primeiro. Por isso as duas colunas
    são coletadas separadas e emparelhadas por posição.
    """
    meses = [(t, x) for t, l, x in caixas(slide) if t >= 0.9 and l < 1.0 and len(x) <= 4]
    textos = [(t, x) for t, l, x in caixas(slide) if t >= 0.9 and l >= 1.0]
    meses.sort(); textos.sort()
    return [[m, txt] for (_, m), (_, txt) in zip(meses, textos)]


# ------------------------------------------------------------------- FOTOS
def _tema_do_slide(slide) -> str:
    """'Obras e melhorias realizadas · Banqueta' -> 'Banqueta'. A legenda fica em
    top~0.63 (mais alta que o corpo dos outros slides, por isso sem filtro de
    top aqui — cortava a única caixa que tem o '·'). Sem separador, usa a
    legenda inteira (raro, mas não pode quebrar por causa disso)."""
    titulo = _titulo(slide)
    for top, left, txt in caixas(slide):
        if txt == titulo:
            continue
        if "·" in txt:
            return txt.split("·", 1)[1].strip()
    return ""


def fotos(prs, indices, mes_tag: str):
    """Extrai as fotos de cada slide de manejo e agrupa por TEMA consecutivo —
    um tema pode ocupar várias slides (ex.: 'Banqueta' em 3 delas na fonte).

    Arquivo sai como '<mes>_fotoNN.jpg': o mês no nome evita um mês novo
    apagar ou colidir com o arquivo do mês anterior (antes era foto01.jpg fixo,
    só servia pra 1 mês por vez)."""
    from PIL import Image
    FOTOS.mkdir(parents=True, exist_ok=True)
    grupos: list[dict] = []
    videos: list[str] = []
    n = 0
    for i in indices:
        slide = prs.slides[i]
        tema = _tema_do_slide(slide)
        if not grupos or grupos[-1]["tema"] != tema:
            grupos.append({"tema": tema, "arquivos": []})
        if any(sh.shape_type is not None and "MEDIA" in str(sh.shape_type) for sh in slide.shapes):
            videos.append(tema)
        for sh in slide.shapes:
            if sh.shape_type is None or "PICTURE" not in str(sh.shape_type):
                continue
            if (sh.width or 0) / EMU_IN < 1.2:      # logo do cabeçalho
                continue
            try:
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
            except Exception:
                continue
            if im.width > LARG_FOTO:
                im = im.resize((LARG_FOTO, round(im.height * LARG_FOTO / im.width)), Image.LANCZOS)
            n += 1
            nome = f"{mes_tag}_foto{n:02d}.jpg"
            im.save(FOTOS / nome, "JPEG", quality=82, optimize=True)
            grupos[-1]["arquivos"].append(f"fotos/{nome}")
    if videos:
        print(f"  [aviso] {len(videos)} tema(s) são VÍDEO, não foto — o pipeline não "
              f"lê vídeo, ficam de fora do deck: {', '.join(videos)}")
    return [g for g in grupos if g["arquivos"]]


def run(pptx_path: Path, mes: str):
    if not pptx_path.exists():
        sys.exit(f"deck não encontrado: {pptx_path}")
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slides = list(prs.slides)

    i_com = _acha(slides, "COMENTÁRIOS — VARIAÇÕES YTD")
    i_prog = _acha(slides, "— PROGRAMAÇÃO")
    is_res = _acha_todos(slides, "RESULTADOS —")
    i_manejo = _acha(slides, "MANEJO — PONTOS DE MELHORIA")
    is_fotos = _acha_todos(slides, "MANEJO — FOTOS E REGISTROS")

    faltando = [nome for nome, v in
                [("comentários YTD", i_com), ("programação de exposições", i_prog),
                 ("manejo", i_manejo)] if v is None]
    if faltando:
        sys.exit(f"não achei o slide de: {', '.join(faltando)} — título mudou no deck?")
    if not is_res:
        print("  [aviso] nenhum slide 'RESULTADOS —' achado — exposições sem resultado este mês?")
    if not is_fotos:
        print("  [aviso] nenhum slide 'MANEJO — FOTOS E REGISTROS' achado")

    dados_mes = {
        "comentarios": comentarios(slides[i_com]),
        "exposicoes": {
            "programacao": programacao(slides[i_prog]),
            "resultados": [resultados(slides[i]) for i in is_res],
        },
        "manejo": manejo(slides[i_manejo]),
        "fotos": fotos(prs, is_fotos, mes.replace("-", "")) if is_fotos else [],
    }

    conteudo = json.loads(OUT.read_text(encoding="utf-8")) if OUT.exists() else {}
    conteudo["_leia"] = (
        "Conteúdo do comitê que NÃO sai de planilha. Semeado do deck oficial de "
        "cada mês por tools/extrair_conteudo.py (roda de novo sem apagar os "
        "outros meses); dá pra editar à mão depois. Cada mês novo é uma chave "
        "'AAAA-MM'; o build usa o mês pedido ou, se não existir, o mês anterior "
        "mais recente. 'fotos' é uma lista de {tema, arquivos}, na ordem do "
        "deck original — tema pode se repetir se o deck usou mais de 1 slide "
        "pra ele.")
    conteudo[mes] = dados_mes
    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text(json.dumps(conteudo, ensure_ascii=False, indent=2), encoding="utf-8")

    n_fotos = sum(len(g["arquivos"]) for g in dados_mes["fotos"])
    print(f"[conteudo] {mes}: {len(dados_mes['comentarios'])} comentários · "
          f"{len(dados_mes['exposicoes']['programacao'])} eventos · "
          f"{len(dados_mes['exposicoes']['resultados'])} exposições com resultado · "
          f"{len(dados_mes['manejo'])} meses de manejo · "
          f"{len(dados_mes['fotos'])} temas de foto ({n_fotos} fotos)")
    print(f"           -> {OUT.relative_to(REPO)}  (fotos em {FOTOS.relative_to(REPO)})")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.exit(__doc__.strip().splitlines()[-1].strip())
    run(Path(sys.argv[1]), sys.argv[2])
