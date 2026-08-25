"""Arquiva o conteúdo dos relatórios mensais já apresentados.

Os decks vivem em `ATA & APRESENTACOES MENSAIS/<ano>/<NN) MES>/` no Drive, e o
Drive não é arquivo morto: pasta é renomeada, arquivo é substituído por uma versão
"final", mês antigo some. Quando isso acontece, o que foi efetivamente apresentado
ao comitê naquele mês deixa de existir — não há de onde reconstruir.

Este script lê os decks e guarda o TEXTO de cada slide, caixa por caixa, com a
posição. Posição importa: nos decks originais do haras o conteúdo não está em
tabela, está em caixas soltas, e sem `top`/`left` não dá para saber o que era
coluna e o que era linha. É o mesmo princípio do arquivo de linhas do fechamento
semanal — guardar a matéria-prima, não só o número que se olhou naquele dia.

O que este script NÃO faz: alimentar o build do comitê. Deck apresentado é registro
histórico, não fonte de dado. A fonte de cada número do comitê é a planilha de
origem, e é ela que a auditoria confere — deck de mês passado como fonte seria
propagar erro antigo para a frente.

Saída: `_cache/decks/<AAAA-MM>_<tipo>.json` (um por deck), cumulativo. Deck já
arquivado não é relido, salvo com `--refazer`.

Uso:
    python tools/arquivar_decks.py                # todos os anos/meses achados
    python tools/arquivar_decks.py 2026           # só um ano
    python tools/arquivar_decks.py 2026 --refazer # reprocessa o que já existe
"""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO / "scripts"))

from _pg_common import DRIVE_ROOT                    # noqa: E402

DECKS_ROOT = DRIVE_ROOT / "ATA & APRESENTACOES MENSAIS"
SAIDA = REPO / "_cache" / "decks"

EMU_IN = 914400

# "04) ABRIL" -> 4. O nome do mês varia (ABRIL, Abril); o número, não.
RE_PASTA_MES = re.compile(r"^(\d{2})\)")

# Só relatório de comitê. A pasta tem ata, DRE, foto de WhatsApp e planilha solta.
RE_DECK = re.compile(r"RELATORIO\s+(MENSAL|TRIMESTRAL)", re.IGNORECASE)


def _caixas(slide) -> list[dict]:
    """Texto de cada caixa, com posição em polegadas, em ordem de leitura."""
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        out.append({"top": round((sh.top or 0) / EMU_IN, 2),
                    "left": round((sh.left or 0) / EMU_IN, 2),
                    "texto": txt})
    return sorted(out, key=lambda c: (c["top"], c["left"]))


def _fotos(slide) -> int:
    return sum(1 for sh in slide.shapes
               if sh.shape_type is not None and "PICTURE" in str(sh.shape_type))


def _extrai(caminho: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(caminho))
    slides = []
    for i, sl in enumerate(prs.slides):
        caixas = _caixas(sl)
        slides.append({
            "n": i,
            # título = a caixa mais alta; nos decks originais é o nome da seção
            # ("EXPOSIÇÕES") e a segunda linha é que diz o assunto ("PROGRAMAÇÃO")
            "titulo": caixas[0]["texto"].split("\n")[0] if caixas else None,
            "fotos": _fotos(sl),
            "caixas": caixas,
        })
    return {
        "arquivo": caminho.name,
        "pasta": caminho.parent.name,
        "modificado": __import__("datetime").datetime.fromtimestamp(
            caminho.stat().st_mtime).isoformat(timespec="seconds"),
        "slides": slides,
    }


def _decks(ano: str | None) -> list[tuple[str, str, Path]]:
    """[(AAAA-MM, tipo, caminho)] dos decks encontrados no Drive."""
    achados = []
    if not DECKS_ROOT.exists():
        sys.exit(f"pasta dos decks não encontrada: {DECKS_ROOT}")
    for pasta_ano in sorted(DECKS_ROOT.iterdir()):
        if not pasta_ano.is_dir() or not re.fullmatch(r"\d{4}", pasta_ano.name):
            continue
        if ano and pasta_ano.name != ano:
            continue
        for pasta_mes in sorted(pasta_ano.iterdir()):
            m = RE_PASTA_MES.match(pasta_mes.name) if pasta_mes.is_dir() else None
            if not m:
                continue
            chave = f"{pasta_ano.name}-{m.group(1)}"
            for arq in sorted(pasta_mes.glob("*.pptx")):
                if arq.name.startswith("~$") or not RE_DECK.search(arq.name):
                    continue
                tipo = "trimestral" if "TRIMESTRAL" in arq.name.upper() else "mensal"
                achados.append((chave, tipo, arq))
    return achados


def run(ano: str | None = None, refazer: bool = False) -> int:
    SAIDA.mkdir(parents=True, exist_ok=True)
    achados = _decks(ano)
    if not achados:
        print(f"[decks] nenhum relatório encontrado em {DECKS_ROOT}" + (f" ({ano})" if ano else ""))
        return 0

    # Mais de um arquivo do mesmo mês e tipo acontece (ex.: junho tem
    # 'TRIMESTRAL' e 'TRIMESTRAL_1'). Fica um por arquivo, com o nome no id —
    # escolher "o certo" seria adivinhar qual foi o apresentado.
    n = 0
    for chave, tipo, arq in achados:
        base = re.sub(r"[^A-Za-z0-9]+", "-", arq.stem).strip("-").lower()
        alvo = SAIDA / f"{chave}_{tipo}_{base}.json"
        if alvo.exists() and not refazer:
            print(f"[decks] {alvo.name} já arquivado")
            continue
        dados = _extrai(arq)
        dados["mes"] = chave
        dados["tipo"] = tipo
        alvo.write_text(json.dumps(dados, ensure_ascii=False, indent=1), encoding="utf-8")
        com_texto = sum(1 for s in dados["slides"] if s["caixas"])
        fotos = sum(s["fotos"] for s in dados["slides"])
        print(f"[decks] {chave} {tipo}: {len(dados['slides'])} slides "
              f"({com_texto} com texto, {fotos} imagens) -> {alvo.name}")
        n += 1
    return n


def main():
    args = [a for a in sys.argv[1:]]
    refazer = "--refazer" in args
    ano = next((a for a in args if re.fullmatch(r"\d{4}", a)), None)
    run(ano, refazer)


if __name__ == "__main__":
    main()
